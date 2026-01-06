from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_visual_premium_presentation_fr():
    prs = Presentation()

    # Define the Premium Color Palette
    COLORS = {
        "NAVY_DEEP": RGBColor(10, 25, 47),
        "GOLD": RGBColor(197, 161, 67),
        "GOLD_LIGHT": RGBColor(230, 190, 100),
        "CRIMSON": RGBColor(180, 0, 40),
        "WHITE": RGBColor(255, 255, 255),
        "BG_LIGHT": RGBColor(248, 250, 252)
    }

    # Image paths from the workspace
    IMAGE_DIR = r"C:\Users\HP\.gemini\antigravity\brain\22708f14-dd02-4f2e-a16a-33482af3c99d"
    
    # Selection of French and African-representative images
    IMAGES = {
        "DASHBOARD": os.path.join(IMAGE_DIR, "executive_dashboard_2026_premium_1767649518351.png"), # Fallback but will update text
        "SUCCESSION": os.path.join(IMAGE_DIR, "succession_pro_fr_senegal_1767652743050.png"),
        "RECOVERY": os.path.join(IMAGE_DIR, "smart_recovery_ia_fr_senegal_1767652719083.png"),
        "WAR_ROOM": os.path.join(IMAGE_DIR, "lexpremium_war_room_tablet_1767649561511.png") # Fallback
    }

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, COLORS["NAVY_DEEP"])
        
        accent_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.1))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = COLORS["GOLD"]
        accent_bar.line.fill.background()

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLORS["GOLD"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(60)
        
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = COLORS["WHITE"]
        subtitle.text_frame.paragraphs[0].font.size = Pt(22)

    def add_visual_content_slide(title_text, points, image_key=None):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, COLORS["BG_LIGHT"])
        
        header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.8))
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS["NAVY_DEEP"]
        header.line.fill.background()

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLORS["GOLD_LIGHT"]
        title.text_frame.paragraphs[0].font.size = Pt(32)
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.word_wrap = True
        
        # Adjust body width if image is present
        if image_key and os.path.exists(IMAGES[image_key]):
            body.width = Inches(4.5)
            body.left = Inches(0.5)
            # Add image
            slide.shapes.add_picture(IMAGES[image_key], Inches(5.2), Inches(1.2), height=Inches(5))
        
        for p in points:
            p_para = tf.add_paragraph()
            p_para.text = "❯ " + p
            p_para.level = 0
            p_para.font.size = Pt(16)
            p_para.font.color.rgb = RGBColor(30, 41, 59)
            p_para.space_after = Pt(8)

    # --- SLIDES GENERATION (STRICT FRENCH) ---
    add_title_slide("LEXPREMIUM 2.0", "L'Intelligence Visuelle au Service du Droit Sénégalais\nInnovation, Performance & Représentation")

    add_visual_content_slide("Tableau de Bord Exécutif", [
        "Pilotage financier 360° en temps réel.",
        "Analyse comparative CA vs Objectifs.",
        "Prévisions de trésorerie par IA Oracle.",
        "Interface Glassmorphism de prestige.",
        "Décisions stratégiques basées sur l'analyse de données."
    ], "DASHBOARD")

    add_visual_content_slide("Recouvrement Intelligent", [
        "Scoring de risque client contextuel (Sénégal).",
        "Interface de relance automatisée multi-canal.",
        "Actions directes : WhatsApp Business & Email.",
        "Réduction drastique des créances impayées.",
        "Monitoring du comportement de paiement en temps réel."
    ], "RECOVERY")

    add_visual_content_slide("Succession Pro Expert", [
        "Moteur de calcul expert (10 méthodes juridiques).",
        "Conformité totale au Code de la Famille sénégalais.",
        "Gestion complexe de la réserve et de l'usufruit.",
        "Génération instantanée de l'acte de partage.",
        "Précision notariale et sécurité juridique garantie."
    ], "SUCCESSION")

    add_visual_content_slide("L'Expérience 'Salle de Commande' (War Room)", [
        "Le palais de justice dans votre poche.",
        "Interface tablette ultra-réactive pour les plaidoiries.",
        "Accès hors-ligne complet pour les tribunaux régionaux.",
        "Dictée vocale et analyse live des pièces adverses.",
        "Standard de prestige international adapté au marché local."
    ], "WAR_ROOM")

    # Save
    prs.save("LexPremium_Visual_Presentation_FR_2026.pptx")
    prs.save("LexPremium_Magistral_PowerPoint_2026.pptx")
    print("French Visual presentations generated successfully with representative images!")

if __name__ == "__main__":
    create_visual_premium_presentation_fr()
