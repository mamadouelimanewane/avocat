from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

def clean_text(text):
    text = text.replace('**', '').replace('*', '')
    text = text.replace('’', "'").replace('“', '"').replace('”', '"').replace('…', '...')
    return re.sub(r'[^\x00-\x7F\xC0-\xFF]', '', text)

def create_pptx(md_file, pptx_file):
    prs = Presentation()
    
    # Color Palette
    navy = RGBColor(15, 23, 42)
    gold = RGBColor(180, 140, 45)
    
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "LEXPREMIUM ERP 2026"
    title.text_frame.paragraphs[0].font.color.rgb = navy
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Guide des Fonctionnalités - L'Excellence Juridique\nPrésentation 360° des Modules"
    
    # Chapters
    chapters = re.split(r'\n## ', content)
    
    for chapter in chapters[1:]:
        lines = chapter.strip().split('\n')
        if not lines: continue
        
        # Slide per Chapter
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title Styling
        title_shape = slide.shapes.title
        chapter_title = clean_text(lines[0])
        title_shape.text = chapter_title.upper()
        
        # Style Title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.bold = True
            paragraph.font.color.rgb = navy
            paragraph.alignment = PP_ALIGN.LEFT

        # Content Styling
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        # Filter content lines
        details = []
        for line in lines[1:]:
            line = line.strip()
            if not line or line.startswith('---') or line.startswith('### '):
                continue
            
            cleaned = clean_text(line)
            
            p = tf.add_paragraph()
            if cleaned.startswith('• ') or cleaned.startswith('- '):
                p.text = cleaned[2:]
                p.level = 1
                p.font.size = Pt(14)
            else:
                p.text = cleaned
                p.level = 0
                p.font.size = Pt(16)
                p.font.italic = True if "Le centre" in cleaned or " module " in cleaned else False
                p.font.color.rgb = RGBColor(80, 80, 80)

    prs.save(pptx_file)
    print(f"PPTX Generated: {pptx_file}")

if __name__ == "__main__":
    create_pptx("c:/gravity/avocat/MANUEL_UTILISATION_DETAILLE.md", "c:/gravity/avocat/release1.pptx")
