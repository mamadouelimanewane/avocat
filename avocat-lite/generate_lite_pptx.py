from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Colors
NAVY = RGBColor(15, 23, 42)
GOLD = RGBColor(180, 140, 45)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(100, 116, 139)
LIGHT_BG = RGBColor(248, 250, 252)

def create_slide(prs, layout_index, title_text, subtitle_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    
    # Set background for all slides to be clean or dark depending on layout
    # Layout 6 is Blank, we usually use it for custom title slides
    
    return slide

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Navy Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background() # No border

    # Diagonal Gold Accent
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(10), 0, Inches(3.333), Inches(3.333))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    shape.rotation = 0

    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "LEXPREMIUM LITE"
    p.font.bold = True
    p.font.size = Pt(72)
    p.font.color.rgb = WHITE
    p.font.name = 'Arial'

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(10), Inches(1))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "L'EXCELLENCE EN TOUTE AUTONOMIE"
    p.font.size = Pt(24)
    p.font.color.rgb = GOLD
    p.font.name = 'Arial'
    
    # Footer/Year
    ft = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(4), Inches(0.5))
    tf = ft.text_frame
    p = tf.paragraphs[0]
    p.text = "ÉDITION 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.font.name = 'Arial'

def add_content_slide(prs, title_str, main_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    
    # Header Bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # Slide Title
    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_str.upper()
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.name = 'Arial'

    # Main Intro Text
    if main_text:
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(1))
        tf = body_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = main_text
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        p.font.name = 'Arial'
        p.font.italic = True

    # Bullet Points
    if bullet_points:
        top_margin = 2.5 if main_text else 1.5
        features_box = slide.shapes.add_textbox(Inches(1), Inches(top_margin), Inches(11), Inches(4.5))
        tf = features_box.text_frame
        tf.word_wrap = True
        
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(51, 65, 85)
            p.font.name = 'Arial'
            p.space_after = Pt(14)
            p.level = 0
            
            # Custom bullet char is hard in python-pptx without workaround, standard bullets will apply

def generate_pptx():
    # Create presentation with widescreen layout (16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title Slide
    add_title_slide(prs)

    # 2. Vision
    add_content_slide(prs, 
        "Vision : Le Cabinet Numérique Duo", 
        "LexPremium Lite est l'édition haute performance conçue pour l'avocat moderne. Un véritable pilote automatique pour votre cabinet.",
        [
            "Centralisez l'intelligence de votre cabinet.",
            "Automatisez la rigueur administrative.",
            "Libérez votre temps pour la stratégie juridique.",
            "Idéal pour l'avocat solo ou accompagné d'un assistant."
        ]
    )

    # 3. Tableau de Bord
    add_content_slide(prs,
        "1. Le Tableau de Bord Intelligents",
        "Un centre de commandement dynamique, pas une simple vue statique.",
        [
            "Focus Actions Requises : Détection des priorités immédiates.",
            "Alertes Forclusions : Surveillance 24/7 des délais de procédure.",
            "KPIs en Temps Réel : Suivi du C.A. et des rentrées de fonds."
        ]
    )

    # 4. Rédaction Assistée
    add_content_slide(prs,
        "2. Rédaction Assistée par IA",
        "Produisez des actes d'une qualité irréprochable en un temps record.",
        [
            "Templates Stratégiques : Assignations et conclusions pré-configurées (OHADA).",
            "Génération Automatique : LexAI rédige le corps des actes selon vos faits.",
            "Export Haute Définition : Documents DOCX prêts à signer."
        ]
    )

    # 5. LexAI
    add_content_slide(prs,
        "3. LexAI : Expertise Augmentée",
        "Une bibliothèque de connaissances juridiques sans précédent.",
        [
            "Jurisprudence Ciblée : Droit sénégalais et Actes Uniformes OHADA.",
            "Analyse de Dossiers : Résumés et extraction de points de droit.",
            "Veille Transversale : Information juridique toujours à jour."
        ]
    )

    # 6. Gestion Dossiers
    add_content_slide(prs,
        "4. Gestion Dossiers & Clients",
        "Une organisation structurée pour une sérénité totale.",
        [
            "Fiches Clients 360° : Historique et coordonnées instantanés.",
            "Suivi de Dossiers : Zéro papier, archivage numérique complet.",
            "Mobilité Totale : Tout votre cabinet accessible sur smartphone."
        ]
    )

    # 7. Facturation
    add_content_slide(prs,
        "5. Facturation & Débours",
        "La rentabilité est le socle de votre indépendance.",
        [
            "Mémoires d'honoraires : Génération en 3 clics.",
            "Gestion des Débours : Suivi rigoureux des frais engagés.",
            "Relances Automatisées : Gestion des impayés par Email/WhatsApp."
        ]
    )

    # 8. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    centered_text = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.333), Inches(2))
    tf = centered_text.text_frame
    p = tf.paragraphs[0]
    p.text = "\"LexPremium : L'excellence technologique au service de votre expertise juridique.\""
    p.font.size = Pt(32)
    p.font.color.rgb = GOLD
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True

    prs.save("c:/gravity/Avocat/avocat-lite/LexPremium_Lite_Presentation_2026.pptx")
    print("PPTX Generated successfully.")

if __name__ == "__main__":
    generate_pptx()
