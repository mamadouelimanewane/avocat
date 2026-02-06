from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_lexpremium_v9_1_presentation():
    prs = Presentation()

    # Define the Premium Color Palette (Deep Navy, Gold, Crimson, Slate)
    COLORS = {
        "NAVY_DEEP": RGBColor(2, 6, 23),      # Slate 950 (LexPremium Dark)
        "GOLD": RGBColor(234, 179, 8),        # Yellow 500
        "GOLD_LIGHT": RGBColor(250, 204, 21), # Yellow 400
        "CRIMSON": RGBColor(220, 38, 38),     # Red 600
        "WHITE": RGBColor(255, 255, 255),
        "SLATE": RGBColor(148, 163, 184),     # Slate 400
        "BG_LIGHT": RGBColor(248, 250, 252)   # Slate 50
    }

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, COLORS["NAVY_DEEP"])
        
        # Add a gold accent bar at the top
        accent_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.12))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = COLORS["GOLD"]
        accent_bar.line.fill.background()

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLORS["GOLD"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(64)
        
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = COLORS["WHITE"]
        subtitle.text_frame.paragraphs[0].font.size = Pt(22)
        
    def add_section_slide(section_name, subtitle=""):
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, RGBColor(10, 10, 10)) # Pure Black for sections
        
        title = slide.shapes.title
        title.text = section_name
        title.text_frame.paragraphs[0].font.color.rgb = COLORS["GOLD"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(54)

        if subtitle:
            sub = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
            p = sub.text_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.italic = True
            p.font.color.rgb = COLORS["WHITE"]
        
        # Add decorative line
        line = slide.shapes.add_shape(1, Inches(1), Inches(4.8), Inches(8), Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS["GOLD"]
        line.line.fill.background()

    def add_content_slide(title_text, subtitle_text, points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, COLORS["BG_LIGHT"])
        
        # Header area
        header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS["NAVY_DEEP"]
        header.line.fill.background()

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLORS["GOLD_LIGHT"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(36)
        
        # Sidebar accent
        sidebar = slide.shapes.add_shape(1, 0, Inches(1.1), Inches(0.18), prs.slide_height - Inches(1.1))
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = COLORS["GOLD"]
        sidebar.line.fill.background()

        # Subtitle textbox
        sub_shape = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(9), Inches(0.4))
        sub_para = sub_shape.text_frame.paragraphs[0]
        sub_para.text = subtitle_text
        sub_para.font.size = Pt(18)
        sub_para.font.bold = True
        sub_para.font.color.rgb = COLORS["NAVY_DEEP"]
        
        body = slide.placeholders[1]
        body.top = Inches(1.8)
        tf = body.text_frame
        tf.word_wrap = True
        
        for p in points:
            p_para = tf.add_paragraph()
            p_para.text = "⚡ " + p
            p_para.level = 0
            p_para.font.size = Pt(17)
            p_para.font.color.rgb = RGBColor(15, 23, 42) # Slate 900
            p_para.space_after = Pt(12)

    # --- SLIDES GENERATION ---
    add_title_slide("LEXPREMIUM PRO v9.1", "L'ARSENAL JURIDIQUE SUPRÊME UNIFIÉ\nDominez le Droit des Affaires International & la Deep Tech")

    add_section_slide("VISION UNIFIÉE 2026", "La Fusion du Savoir Cabinet et de l'Index Mondial")
    
    add_content_slide("L'Arsenal des 10,000+ Actes", "Une Couverture Industrielle Sans Précédent", [
        "Index Suprême : 33 pôles économiques indexés (Deep Tech, BioTech, Énergie, Maritime).",
        "Standards Mondiaux : Accès immédiat aux contrats de niveau international (PPA, Charter Party).",
        "Bibliothèque de Prestiges : De la constitution de startup à l'acquisition minière.",
        "Unification v9.1 : Vos propres modèles fusionnés dans l'index de recherche global.",
        "Élite Opérationnelle : Réduisez le temps de rédaction de 90% sur les dossiers complexes."
    ])

    add_section_slide("I. SENTINELLE CRAWLER™", "La Veille Juridique Assistée par IA")
    add_content_slide("Sentinelle v9.0 : Le Pouvoir de l'Anticipation", "Ne Suivez Plus le Droit, Devancez-le", [
        "Crawler Automatisé : Scan quotidien des journaux officiels et flux mondiaux.",
        "Injection Live : Nouveaux standards contractuels ajoutés sans intervention humaine.",
        "Analyse Comparative : Détection des changements de clauses types sur le marché.",
        "Alerte Conformité : Notification immédiate en cas de mise à jour réglementaire.",
        "Intelligence Collective : Bénéficiez des actifs juridiques les plus pointus au monde."
    ])

    add_section_slide("II. IA JURIDIQUE & LEXAI", "L'Avocat Augmenté par l'IA Cognitive")
    add_content_slide("LexAI Predict & Scan Adverse", "Science de la Victoire & Renseignement", [
        "LexAI Predict : Calcul des chances de succès basé sur l'historique jurisprudentiel.",
        "Scanner Adverse : Extraction instantanée des points faibles des conclusions adverses.",
        "Rédaction Augmentée : Complétion intelligente d'actes basée sur votre style.",
        "Sherlock OSINT : Renseignement d'affaires sur la solvabilité de la partie adverse.",
        "Nexus Graph : Visualisation des réseaux d'influence et conflits d'intérêts."
    ])

    add_section_slide("III. GESTION EXPERTE & MÉTIER", "Moteurs de Calcul de Haute Précision")
    add_content_slide("Succession Pro™ & Quantum", "La Rigueur des Algorithmes", [
        "Succession Pro : Liquidation successorale complexe (10 méthodes) - Gain de temps 95%.",
        "Moteur OHADA : Calculs financiers et intérêts légaux automatisés.",
        "Quantum Simulator : Simulations d'indemnités sociales et foncières ultra-précises.",
        "Conformité Fiscale : Barèmes CGI et taxes transactionnelles intégrés.",
        "Rapports d'Expertise : Documents PDF d'élite prêts pour le dépôt au greffe."
    ])

    add_section_slide("IV. FINANCE & RECOUVREMENT", "Ingénierie de Trésorerie & Mobilité")
    add_content_slide("Recovery Engine & SYSCOHADA", "L'Efficacité Transformée en Cash", [
        "Smart Recovery™ : Relances WhatsApp/Email automatisées et scoring risque client.",
        "Execution Commander : Carte interactive des actifs saisis et pilotage des huissiers.",
        "Cashflow Oracle : Prédiction des entrées de fonds à 90 jours via IA.",
        "Snap-to-SYSCOHADA : Intégration comptable instantanée par OCR de factures.",
        "War Room Tablette : Votre cabinet en poche pour les audiences hors-ligne."
    ])

    add_section_slide("V. EXPÉRIENCE CLIENT ELITE", "Le Portail VIP des Grands Mandants")
    add_content_slide("Souveraineté & Transparence", "Le Luxe Logiciel pour vos Clients", [
        "Procedure Tracker : Suivi 'Amazon' de l'avancement pour le client.",
        "Extranet Sécurisé : Espace VIP pour le dépôt de pièces et suivi financier.",
        "Centif Guard : Protection anti-blanchiment (AML) intégrée sur les fonds CARPA.",
        "WhatsApp Pro Bridge : Notifications de procédure sur le canal favori des clients.",
        "Chiffrement AES-256 : Sécurité de grade militaire pour le secret professionnel."
    ])

    add_title_slide("DOMINEZ VOTRE AVENIR", "LexPremium v9.1\nL'Excellence est une Décision Stratégique.\nDémo : avocat-tito.vercel.app")

    # Save file
    output_filename = "LexPremium_Unification_Magistrale_2026.pptx"
    prs.save(output_filename)
    print(f"Unified v9.1 Presentation Generated: {output_filename}")

if __name__ == "__main__":
    create_lexpremium_v9_1_presentation()
