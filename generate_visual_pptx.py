from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_visual_premium_presentation():
    prs = Presentation()

    # Define the Premium Color Palette
    COLORS = {
        "NAVY_DEEP": RGBColor(10, 25, 47),
        "GOLD": RGBColor(197, 161, 67),
        "GOLD_LIGHT": RGBColor(230, 190, 100),
        "CRIMSON": RGBColor(180, 0, 40),
        "WHITE": RGBColor(255, 255, 255),
        "BG_LIGHT": RGBColor(248, 250, 252)
    }

    # Image paths from the workspace
    IMAGE_DIR = r"C:\Users\HP\.gemini\antigravity\brain\22708f14-dd02-4f2e-a16a-33482af3c99d"
    IMAGES = {
        "DASHBOARD": os.path.join(IMAGE_DIR, "executive_dashboard_2026_premium_1767649518351.png"),
        "SUCCESSION": os.path.join(IMAGE_DIR, "succession_pro_expert_tool_1767649544616.png"),
        "RECOVERY": os.path.join(IMAGE_DIR, "smart_recovery_ia_2_0_ui_1767649585253.png"),
        "WAR_ROOM": os.path.join(IMAGE_DIR, "lexpremium_war_room_tablet_1767649561511.png")
    }

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, COLORS["NAVY_DEEP"])
        
        accent_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.1))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = COLORS["GOLD"]
        accent_bar.line.fill.background()

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLORS["GOLD"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(60)
        
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = COLORS["WHITE"]
        subtitle.text_frame.paragraphs[0].font.size = Pt(22)

    def add_visual_content_slide(title_text, points, image_key=None):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, COLORS["BG_LIGHT"])
        
        header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.8))
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS["NAVY_DEEP"]
        header.line.fill.background()

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLORS["GOLD_LIGHT"]
        title.text_frame.paragraphs[0].font.size = Pt(32)
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.word_wrap = True
        
        # Adjust body width if image is present
        if image_key and os.path.exists(IMAGES[image_key]):
            body.width = Inches(4.5)
            body.left = Inches(0.5)
            # Add image
            slide.shapes.add_picture(IMAGES[image_key], Inches(5.2), Inches(1.2), height=Inches(5))
        
        for p in points:
            p_para = tf.add_paragraph()
            p_para.text = "❯ " + p
            p_para.level = 0
            p_para.font.size = Pt(16)
            p_para.font.color.rgb = RGBColor(30, 41, 59)
            p_para.space_after = Pt(8)

    # --- SLIDES GENERATION ---
    add_title_slide("LEXPREMIUM 2.0", "L'Intelligence Visuelle au Service du Droit\nGuide de Démonstration Magistrale")

    add_visual_content_slide("Executive Dashboard V2", [
        "Pilotage financier 360°.",
        "Analyse Revenue vs Objectifs.",
        "Prévisions IA glowing analytics.",
        "Interface Glassmorphism Moderne.",
        "Décisions stratégiques basées data."
    ], "DASHBOARD")

    add_visual_content_slide("Smart Recovery Engine", [
        "Scoring de risque IA (0-100).",
        "Interface de relance automatisée.",
        "Boutons d'action WhatsApp & Email.",
        "Réduction drastique des impayés.",
        "Monitoring comportement payeur."
    ], "RECOVERY")

    add_visual_content_slide("Succession Pro Expert", [
        "Calculateur expert 10 méthodes.",
        "Visualisation arbre généalogique.",
        "Partage complexe réserve/usufruit.",
        "Génération acte de partage Gold.",
        "Précision notariale garantie."
    ], "SUCCESSION")

    add_visual_content_slide("L'Expérience War Room", [
        "Le palais de justice dans la poche.",
        "Interface tablette ultra-réactive.",
        "Mode offline pour les plaidoiries.",
        "Dictée vocale et analyse live.",
        "Standard de prestige international."
    ], "WAR_ROOM")

    # Save
    prs.save("LexPremium_Visual_Presentation_2026.pptx")
    prs.save("LexPremium_Magistral_PowerPoint_2026.pptx") # Overwrite with visual version
    print("Visual presentations generated with success!")

if __name__ == "__main__":
    create_visual_premium_presentation()
