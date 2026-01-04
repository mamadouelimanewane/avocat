from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_ultimate_magistral_presentation():
    prs = Presentation()

    # Define the TITO Website Color Palette (Gold, Black, Red)
    COLORS = {
        "BLACK": RGBColor(15, 23, 42),       # Slate 900
        "PURE_BLACK": RGBColor(0, 0, 0),
        "GOLD": RGBColor(184, 134, 11),      # Dark Goldenrod
        "GOLD_BRIGHT": RGBColor(218, 165, 32),
        "RED": RGBColor(220, 38, 38),        # Red
        "WHITE": RGBColor(255, 255, 255),
        "SLATE_BG": RGBColor(248, 250, 252)
    }

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_slide(title_text, subtitle_text, color=COLORS["BLACK"]):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, color)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLORS["GOLD_BRIGHT"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(64)
        
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = COLORS["WHITE"]
        subtitle.text_frame.paragraphs[0].font.size = Pt(28)
        
    def add_section_slide(section_name, color=COLORS["PURE_BLACK"]):
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, color)
        
        title = slide.shapes.title
        title.text = section_name
        title.text_frame.paragraphs[0].font.color.rgb = COLORS["GOLD"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(50)
        
        line = slide.shapes.add_shape(1, Inches(1), Inches(4.5), Inches(8), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS["RED"]

    def add_content_slide(title_text, subtitle_text, points, accent_color=COLORS["GOLD"]):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, COLORS["SLATE_BG"])
        
        header_strip = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.8))
        header_strip.fill.solid()
        header_strip.fill.fore_color.rgb = COLORS["BLACK"]
        header_strip.line.fill.background()

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLORS["GOLD_BRIGHT"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(36)
        
        sidebar = slide.shapes.add_shape(1, 0, Inches(0.8), Inches(0.2), prs.slide_height - Inches(0.8))
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = COLORS["RED"]
        sidebar.line.fill.background()

        sub_shape = slide.shapes.add_textbox(Inches(0.6), Inches(1), Inches(9), Inches(0.5))
        sub_para = sub_shape.text_frame.paragraphs[0]
        sub_para.text = subtitle_text
        sub_para.font.size = Pt(18)
        sub_para.font.bold = True
        sub_para.font.color.rgb = COLORS["BLACK"]
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.word_wrap = True
        
        for p in points:
            p_para = tf.add_paragraph()
            p_para.text = "◈ " + p
            p_para.level = 0
            p_para.font.size = Pt(19)
            p_para.font.color.rgb = COLORS["BLACK"]
            p_para.space_after = Pt(12)

    # --- SLIDES GENERATION ---
    add_title_slide("LEXPREMIUM", "ERP Juridique de Niveau Magistral\nL'Art de la Justice Propulsé par l'IA", COLORS["PURE_BLACK"])

    add_section_slide("ACCUEIL & VISION")
    add_content_slide("L'Écosystème LexPremium", "Une Vision à 360° du Cabinet d'Elite", [
        "Gestion Intégrée : Dossiers, Clients, Audiences.",
        "IA Stratégique : LexAI, Scan Adverse, Aide à la rédaction.",
        "Documentation Élite : GED chiffrée & Bible juridique.",
        "Communication & Portail : WhatsApp, Extranet, Newsletters.",
        "Ressources Humaines : Talents, Utilisateurs & Finance OHADA."
    ])

    add_section_slide("I. DOCUMENTATION & GED ÉLITE")
    add_content_slide("GED Élite & Archivage", "Maîtrise Documentaire Totale", [
        "OCR Haute-Résolution : Transformation immédiate des scans en textes vivants.",
        "Indexation Sémantique : Retrouvez n'importe quelle clause ou pièce par concept.",
        "Versioning Dynamique : Suivez chaque modification d'acte (v1, v2, Final).",
        "Coffre-Fort Numérique : Archivage chiffré (AES-256) des documents sensibles."
    ])
    add_content_slide("Bible Juridique & Modèles", "Standardisez l'Excellence du Cabinet", [
        "Base de Modèles d'Élite : Centralisez les contrats-types et mémoires.",
        "Génération Dynamique : Liaison automatique client <-> dossier <-> acte.",
        "Aide à la Rédaction IA : 75% du travail de structure fait par l'IA en 3 clics.",
        "Bibliothèque de Clauses : Insérez des phrases sécurisées et testées légalement."
    ])

    add_section_slide("II. COMMUNICATION & PORTAIL CLIENT")
    add_content_slide("Communication Multi-Canal", "Réactivité et Traçabilité Sans Faille", [
        "Messagerie Unifiée : Consolidation des Emails et WhatsApp par dossier.",
        "Notifications Automatisées : Rappels Audiences et RDV par SMS/Email.",
        "Newsletters Juridiques : Informez vos clients des veilles législatives.",
        "Collaboration Interne : Messagerie sécurisée entre collaborateurs."
    ])
    add_content_slide("Portail Client Extra-Premium", "Le Luxe de la Transparence", [
        "Extranet Sécurisé : Vos clients suivent leur procédure en temps réel 24h/24.",
        "Dépôt de Pièces : Le mandant télécharge ses documents directement en GED.",
        "Messagerie Chiffrée : Échanges confidentiels sanctuarisés hors emails.",
        "Suivi Financier : Consultation des provisions et factures par le mandant."
    ])

    add_section_slide("III. INTELLIGENCE STRATÉGIQUE & IA DISRUPTIVE")
    add_content_slide("LexAI Predict & Magistrat Intel", "La Science de la Victoire", [
        "LexAI Predict : Probabilité de succès basée sur la jurisprudence OHADA/Sénégal.",
        "Magistrat Intel : Profilage psychométrique des juges (Sévérité, Équité, Rapidité).",
        "Stratégie de Plaidoirie : Conseils personnalisés selon le profil du magistrat.",
        "Analyse de Risque 360° : Identification des failles dans l'argumentaire adverse."
    ])
    add_content_slide("Sherlock OSINT & Nexus Graph", "Renseignement & Réseaux d'Influence", [
        "Sherlock OSINT : Scan automatisé du patrimoine et du web gris de la partie adverse.",
        "Nexus Graph : Visualisation interactive des liens familiaux, d'affaires et d'influence.",
        "Détection des Conflits d'Intérêts : Alerte immédiate sur les liens cachés.",
        "Solvabilité Adverse : Localisation d'actifs pour les saisies conservatoires."
    ])
    add_content_slide("Quantum Simulator & War Room", "Domination en Négociation et à l'Audience", [
        "Quantum Simulator : Projection ultra-visuelle des indemnités et dommages-intérêts.",
        "Simulation de Scénarios : Impact instantané du changement de paramètres légaux.",
        "War Room Tablette : Interface immersive optimisée pour les plaidoiries au Palais.",
        "Voice Commander : Pilotage vocal du dossier en situation de mobilité."
    ])

    add_section_slide("IV. FINANCE, RÉCUPÉRATION & CONFORMITÉ")
    add_content_slide("Execution Commander", "Le Pilotage de l'Exécution Tactique", [
        "Carte Tactique de Dakar : Visualisation géographique des cibles de saisie.",
        "Suivi des Huissiers : Statut live des interventions sur le terrain.",
        "Compteur de Recouvrement : Taux de transformation du jugement en cash.",
        "Priorisation par IA : Focus automatique sur les actifs les plus solvables."
    ])
    add_content_slide("Centif Guard & Cashflow Oracle", "Sécurité Pénale et Vision Financière", [
        "Centif Guard : Scanner AML intégré pour chaque virement de fonds CARPA.",
        "Conformité Instantanée : Vérification contre les listes de sanctions et PEPs.",
        "Cashflow Oracle : Prédiction de trésorerie à 6 mois (IA comportementale).",
        "Snap-to-SYSCOHADA : OCR comptable intelligent avec imputation automatique."
    ])

    add_section_slide("V. EXPÉRIENCE CLIENT & AUTOMATISATION")
    add_content_slide("Procedure Tracker & WhatsApp Bridge", "Transparence et Réactivité Totale", [
        "Procedure Tracker : Suivi style 'Amazon' de l'avancement pour le client.",
        "WhatsApp Pro Bridge : Communication instantanée via templates pré-rédigés.",
        "Parapheur Numérique : Circuit de validation documentaire (Rédaction -> Signature).",
        "Auto-Classifier GED : Classement sémantique automatique des fichiers déposés."
    ])

    add_section_slide("VI. UTILISATEURS, RH & INFRASTRUCTURE")
    add_content_slide("Gestion des Utilisateurs", "Sécurité et Hiérarchie du Cabinet", [
        "Rôles & Permissions : Gérez les accès (Associés, Avocats, Secrétaires).",
        "Activity Log : Traçabilité complète des actions par utilisateur.",
        "Gestion MFA : Authentification multifacteur de grade bancaire.",
        "Audit Trail : Historique inaltérable pour le respect de la déontologie."
    ])
    add_content_slide("Gestion RH & Finance", "Rigueur OHADA & Talents", [
        "Comptabilité SYSCOHADA : Écritures automatiques et gestion CARPA.",
        "Gestion RH : Performance, congés et rentabilité par collaborateur.",
        "Time-Tracking : Capture de chaque minute valorisable au dossier.",
        "Facturation Prestige : Devis et factures au standing de votre marque."
    ])

    add_section_slide("VII. PARAMÉTRAGES & AUDIT")
    add_content_slide("Administration Centrale", "Le Cabinet sur Mesure", [
        "Branding Cabinet : Personnalisation logo, charte et entêtes.",
        "Audit Trail Complet : Historique inaltérable pour la déontologie.",
        "Monitoring Système : Santé des données et backups en temps réel.",
        "Sécurisation Totale : Paramétrages avancés de la souveraineté."
    ])

    add_title_slide("DOMINEZ VOTRE AVENIR", "LexPremium 2026\nL'Excellence est une Décision.\nDémo : avocat-tito.vercel.app", COLORS["PURE_BLACK"])

    # Final Output File
    output_file = "LexPremium_Magistral_PowerPoint_2026.pptx"
    prs.save(output_file)
    print(f"Presentation updated successfully: {output_file}")

if __name__ == "__main__":
    create_ultimate_magistral_presentation()
