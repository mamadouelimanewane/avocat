from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import re
import os

# --- STYLE CONSTANTS ---
NAVY = RGBColor(10, 37, 64)      # #0A2540
GOLD = RGBColor(212, 175, 55)    # #D4AF37
TECH_BLUE = RGBColor(0, 163, 224) # #00A3E0
TEXT_DARK = RGBColor(31, 41, 55) # #1F2937
BG_COLOR = RGBColor(248, 250, 252) # #F8FAFC
BROWN = RGBColor(101, 67, 33)

def clean_text(text):
    text = text.replace('**', '').replace('*', '')
    text = text.replace('’', "'").replace('“', '"').replace('”', '"').replace('…', '...')
    return re.sub(r'[^\x00-\x7F\xC0-\xFF]', '', text)

def add_background_rect(slide, color):
    # Use a large rectangle as background for stability
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background() # No border

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    add_background_rect(slide, NAVY)
    
    # Gold strip
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = GOLD
    rect.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "LEXPREMIUM ERP 2026"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "L'Excellence Juridique a l'Ere du Numerique"
    p.font.size = Pt(32)
    p.font.italic = True
    p.font.color.rgb = GOLD
    
    # Year
    year_box = slide.shapes.add_textbox(Inches(7), Inches(6.5), Inches(2.5), Inches(0.5))
    tf = year_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Version Strategique 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.RIGHT

def add_module_slide(prs, num, title, catchphrase, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    add_background_rect(slide, BG_COLOR)
    
    # Header Band (Navy)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    
    # Number Circle (Gold)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), Inches(0.2), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    
    num_txt = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(0.8), Inches(0.8))
    tf = num_txt.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE # Corrected enum
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Module Title
    title_txt = slide.shapes.add_textbox(Inches(1.3), Inches(0.3), Inches(8), Inches(0.6))
    tf = title_txt.text_frame
    p = tf.paragraphs[0]
    p.text = clean_text(title).upper()
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Catchphrase (Tech Blue)
    catch_txt = slide.shapes.add_textbox(Inches(1.3), Inches(0.85), Inches(8), Inches(0.4))
    tf = catch_txt.text_frame
    p = tf.paragraphs[0]
    p.text = clean_text(catchphrase)
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = TECH_BLUE
    
    # Left Column: Image insertion if it exists
    img_path = f"c:/gravity/avocat/icons/module_{num}.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.7), Inches(2.2), width=Inches(2.6))

    # Right Column: Content Box
    right_x = Inches(4)
    right_y = Inches(1.8)
    right_w = Inches(5.5)
    
    content_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, right_y, right_w, Inches(4.5))
    content_rect.fill.solid()
    content_rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_rect.line.color.rgb = GOLD
    content_rect.line.width = Pt(1)
    
    content_box = slide.shapes.add_textbox(right_x + Inches(0.2), right_y + Inches(0.1), right_w - Inches(0.4), Inches(4.3))
    content_box.text_frame.word_wrap = True
    
    for i, item in enumerate(items[:8]): 
        p = content_box.text_frame.add_paragraph()
        p.space_before = Pt(12)
        
        cleaned_item = clean_text(item)
        if ':' in cleaned_item:
            label, desc = cleaned_item.split(':', 1)
            run = p.add_run()
            run.text = "• " + label.strip() + " : "
            run.font.bold = True
            run.font.color.rgb = BROWN
            run.font.size = Pt(18)
            
            run2 = p.add_run()
            run2.text = desc.strip()
            run2.font.size = Pt(16)
            run2.font.color.rgb = TEXT_DARK
        else:
            p.text = "• " + cleaned_item
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_DARK

def create_premium_pptx(md_file, pptx_file):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    add_title_slide(prs)
    
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    chapters = re.split(r'\n## ', content)
    
    for chapter in chapters[1:]:
        lines = [l.strip() for l in chapter.strip().split('\n') if l.strip()]
        if not lines: continue
        
        header = lines[0]
        match = re.match(r'(\d+)\.\s+(.*)', header)
        if match:
            num = match.group(1)
            title = match.group(2)
        else:
            num = "?"
            title = header
            
        catchphrase = ""
        items = []
        for line in lines[1:]:
            if line.startswith('### '): continue
            if line.startswith('• ') or line.startswith('- '):
                items.append(line[2:])
            elif not catchphrase:
                catchphrase = line
            else:
                items.append(line)
        
        add_module_slide(prs, num, title, catchphrase, items)
        
    prs.save(pptx_file)
    print(f"Premium PPTX Generated: {pptx_file}")

if __name__ == "__main__":
    create_premium_pptx("c:/gravity/avocat/MANUEL_UTILISATION_DETAILLE.md", "c:/gravity/avocat/release1_premium.pptx")
