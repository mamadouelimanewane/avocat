from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_ultimate_magistral_presentation_v2():
    prs = Presentation()

    # Define the Premium Color Palette (Deep Navy, Gold, Crimson, Slate)
    COLORS = {
        "NAVY_DEEP": RGBColor(10, 25, 47),    # Ultra Deep Navy
        "GOLD": RGBColor(197, 161, 67),      # Metallic Gold
        "GOLD_LIGHT": RGBColor(230, 190, 100),
        "CRIMSON": RGBColor(180, 0, 40),     # Professional Red
        "WHITE": RGBColor(255, 255, 255),
        "SLATE": RGBColor(71, 85, 105),
        "BG_LIGHT": RGBColor(248, 250, 252)
    }

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, COLORS["NAVY_DEEP"])
        
        # Add a gold accent bar at the top
        accent_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.1))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = COLORS["GOLD"]
        accent_bar.line.fill.background()

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLORS["GOLD"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(64)
        
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = COLORS["WHITE"]
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        
    def add_section_slide(section_name):
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, COLORS["NAVY_DEEP"])
        
        title = slide.shapes.title
        title.text = section_name
        title.text_frame.paragraphs[0].font.color.rgb = COLORS["GOLD"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(54)
        
        # Add decorative line
        line = slide.shapes.add_shape(1, Inches(1), Inches(4.5), Inches(8), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS["CRIMSON"]
        line.line.fill.background()

    def add_content_slide(title_text, subtitle_text, points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, COLORS["BG_LIGHT"])
        
        # Header area
        header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS["NAVY_DEEP"]
        header.line.fill.background()

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLORS["GOLD_LIGHT"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(36)
        
        # Sidebar accent
        sidebar = slide.shapes.add_shape(1, 0, Inches(1.0), Inches(0.15), prs.slide_height - Inches(1.0))
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = COLORS["CRIMSON"]
        sidebar.line.fill.background()

        # Subtitle textbox
        sub_shape = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(9), Inches(0.4))
        sub_para = sub_shape.text_frame.paragraphs[0]
        sub_para.text = subtitle_text
        sub_para.font.size = Pt(18)
        sub_para.font.bold = True
        sub_para.font.color.rgb = COLORS["NAVY_DEEP"]
        
        body = slide.placeholders[1]
        body.top = Inches(1.6)
        tf = body.text_frame
        tf.word_wrap = True
        
        for p in points:
            p_para = tf.add_paragraph()
            p_para.text = "❯ " + p
            p_para.level = 0
            p_para.font.size = Pt(18)
            p_para.font.color.rgb = RGBColor(30, 41, 59) # Slate 800
            p_para.space_after = Pt(10)

    # --- SLIDES GENERATION ---
    add_title_slide("LEXPREMIUM 2.0", "L'EXCELLENCE OPÉRATIONNELLE AUGMENTÉE\nTableau de Bord Exécutif | Recouvrement IA | Succession Pro")

    add_section_slide("ACCUEIL & VISION")
    add_content_slide("Le Centre de Commandement", "Une Vision à 360° du Cabinet en 2026", [
        "Intelligence Artificielle Intégrée : Analyse de pièces, rédaction, prédictions.",
        "Gestion Financière de Pointe : Pilotage par la donnée (Data-Driven).",
        "Expertise Juridique Automatisée : Moteurs de calcul successoral conformes.",
        "Expérience Client Ultra-Premium : Transparence totale via le portail VIP.",
        "Souveraineté & Sécurité : Cloud souverain et chiffrement AES-256."
    ])

    add_section_slide("I. TABLEAU DE BORD EXÉCUTIF (V2.0)")
    add_content_slide("Pilotage Stratégique", "Le Cockpit Décisionnel des Associés", [
        "KPI Temps Réel : Chiffre d'Affaires réalisé vs Objectifs fixés.",
        "Marge Nette de Dossier : Rentabilité réelle calculée sur le coût collaborateur.",
        "Trésorerie Oracle : Prévisions d'encaissements à 30, 60 et 90 jours par IA.",
        "Scoring de Risque Global : Identification proactive des faiblesses du cabinet.",
        "Alertes Intelligentes : Notifications sur les créances critiques et retards."
    ])

    add_section_slide("II. SMART RECOVERY ENGINE™ (V2.0)")
    add_content_slide("Transformer les Créances en Cash", "Réduisez vos Impayés de 60%", [
        "Scoring Client IA (0-100) : Analyse comportementale du payeur.",
        "Relances Multi-Canaux : Séquences automatisées via WhatsApp Business & Email.",
        "Mise en Demeure Instantanée : Génération du document juridique en 1 clic.",
        "Priorisation Tactique : L'IA vous indique qui relancer et quand pour maximiser le cash.",
        "Traçabilité Totale : Historique des reliances directement lié au dossier."
    ])

    add_section_slide("III. SUCCESSION PRO™ (V2.0)")
    add_content_slide("L'Expertise Liquidation", "De 3 jours à 9 minutes par Dossier", [
        "Moteur 10 Méthodes : Liquidation de régime, réserve, usufruit, libéralités.",
        "Conformité Code de la Famille : Algorithmes validés juridiquement.",
        "Calcul Fiscal Automatique : Droits de succession calculés selon le barème CGI.",
        "Arbre des Héritiers Intelligent : Gestion des ordres et des représentations.",
        "Export Acte de Partage : Rapport professionnel prêt pour signature."
    ])

    add_section_slide("IV. INTELLIGENCE ARTIFICIELLE & GED")
    add_content_slide("LexAI & OCR Neural", "Votre Associé de Génie 24h/24", [
        "OCR Neural Haute Précision : Extraction full-text des factures et pièces.",
        "Analyse de Conclusions Adverses : Détection de failles et contradictions.",
        "JusticePredictor™ : Simulation des chances de succès basée sur la data.",
        "Contract Analyzer : Vérification de conformité OHADA systématique.",
        "Cache Sémantique : Réponses IA instantanées et réduction des coûts de 80%."
    ])

    add_section_slide("V. EXPÉRIENCE CLIENT & MOBILITÉ")
    add_content_slide("Relation Client VIP", "Le Luxe au Service de la Transparence", [
        "Procedure Tracker : Suivi d'avancement Amazon-style pour le mandant.",
        "WhatsApp Bridge sécurisé : Échanges confidentiels historisés au dossier.",
        "War Room Tablette : Accès offline complet pour les plaidoiries au Palais.",
        "Snap-to-SYSCOHADA : Numérisation des frais et auto-imputation comptable.",
        "Portail Client 5* : Dépôt sécurisé, paiements en ligne, suivi financier."
    ])

    add_section_slide("VI. SÉCURITÉ & CONFORMITÉ")
    add_content_slide("Le Socle de Confiance", "La Déontologie au Coeur de la Tech", [
        "Centif Guard™ : Monitoring anti-blanchiment pour les fonds CARPA.",
        "Audit Trail Inaltérable : Traçabilité historique de chaque accès.",
        "Chiffrement AES-256 : Protection de grade militaire pour vos données.",
        "Hébergement Souverain : Données résidant sur des infrastructures sécurisées.",
        "MFA (Authentification Multifacteur) : Protection contre l'usurpation d'identité."
    ])

    add_title_slide("REJOIGNEZ L'ELITE", "LexPremium 2.0\nLe Futur de la Profession est Déjà Là.\nContact : www.lexpremium.sn")

    # Save files
    output_file = "LexPremium_Presentation.pptx"
    output_magistral = "LexPremium_Magistral_PowerPoint_2026.pptx"
    prs.save(output_file)
    prs.save(output_magistral)
    print(f"Presentations updated: {output_file} and {output_magistral}")

if __name__ == "__main__":
    create_ultimate_magistral_presentation_v2()
